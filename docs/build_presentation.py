"""Generate the PRP presentation deck.

10-minute talk for teachers + peers. Slides use the same dark/indigo/purple
palette as the React frontend (tokens.css). Plain English — no jargon walls.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ----------------------------------------------------------------------------
# Theme — matches frontend/src/theme/tokens.css
# ----------------------------------------------------------------------------
BG_DEEP        = RGBColor(0x05, 0x06, 0x08)
BG             = RGBColor(0x0a, 0x0d, 0x14)
SURFACE        = RGBColor(0x14, 0x1a, 0x26)
BORDER         = RGBColor(0x33, 0x33, 0x44)
TEXT           = RGBColor(0xed, 0xed, 0xf2)
TEXT_MUTED     = RGBColor(0x8a, 0x8f, 0x98)
ACCENT_PURPLE  = RGBColor(0x94, 0x57, 0xff)
ACCENT_INDIGO  = RGBColor(0x63, 0x66, 0xf1)

SNAKE_BLUE     = RGBColor(0x1e, 0x3a, 0x8a)
SNAKE_RED      = RGBColor(0xef, 0x44, 0x44)
SNAKE_WHITE    = RGBColor(0xff, 0xff, 0xff)

SEV_INFO       = RGBColor(0x10, 0xb9, 0x81)
SEV_HIGH       = RGBColor(0xf9, 0x73, 0x16)
SEV_CRIT       = RGBColor(0xef, 0x44, 0x44)
SEV_MED        = RGBColor(0xea, 0xb3, 0x08)
SEV_LOW        = RGBColor(0x3b, 0x82, 0xf6)

# Plotly-style category palette for the "What we detect" chips
CAT_COLORS = {
    "BENIGN":       RGBColor(0x10, 0xb9, 0x81),
    "DDoS":         RGBColor(0xef, 0x44, 0x44),
    "DoS":          RGBColor(0xf9, 0x73, 0x16),
    "PortScan":     RGBColor(0x06, 0xb6, 0xd4),
    "BruteForce":   RGBColor(0xf5, 0x9e, 0x0b),
    "WebAttack":    RGBColor(0xea, 0xb3, 0x08),
    "Botnet":       RGBColor(0xa8, 0x55, 0xf7),
    "Infiltration": RGBColor(0xec, 0x48, 0x99),
}

FONT_BODY = "Inter"
FONT_FALLBACK = "Calibri"   # PowerPoint uses this if Inter isn't installed
FONT_MONO = "JetBrains Mono"


# ----------------------------------------------------------------------------
# Slide builders
# ----------------------------------------------------------------------------
def add_dark_bg(slide, prs):
    """Solid dark background + a soft purple corner glow + a thin accent bar."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False

    # Purple corner glow (top-left) — an oversized soft circle clipped by slide
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-3.5), Inches(-3), Inches(6), Inches(6)
    )
    glow.line.fill.background()
    glow.fill.solid()
    glow.fill.fore_color.rgb = ACCENT_PURPLE
    glow.fill.transparency = 0  # python-pptx limitation; set via xml below
    _set_shape_alpha(glow, 0.12)

    # Indigo corner glow (bottom-right)
    glow2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        prs.slide_width - Inches(3),
        prs.slide_height - Inches(3),
        Inches(6),
        Inches(6),
    )
    glow2.line.fill.background()
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = ACCENT_INDIGO
    _set_shape_alpha(glow2, 0.08)

    # Thin accent strip across the top (the same purple bar under the navbar)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.07)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_PURPLE


def _set_shape_alpha(shape, alpha: float) -> None:
    """Set fill alpha (python-pptx doesn't expose this directly)."""
    from pptx.oxml.ns import qn
    sp = shape.fill._xPr
    fill = sp.find(qn("a:solidFill"))
    if fill is None:
        return
    color = fill.find(qn("a:srgbClr"))
    if color is None:
        return
    # alpha is per-mille (0–100000)
    alpha_node = color.find(qn("a:alpha"))
    if alpha_node is None:
        from lxml import etree
        alpha_node = etree.SubElement(color, qn("a:alpha"))
    alpha_node.set("val", str(int(alpha * 100000)))


def add_text(
    slide,
    left, top, width, height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    left, top, width, height,
    items: list[str],
    *,
    size: int = 18,
    color: RGBColor = TEXT,
    bullet_color: RGBColor = ACCENT_PURPLE,
    line_spacing: float = 1.35,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        # bullet glyph
        b = p.add_run()
        b.text = "●  "
        b.font.name = FONT_BODY
        b.font.size = Pt(size)
        b.font.color.rgb = bullet_color
        # body
        r = p.add_run()
        r.text = item
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_card(
    slide,
    left, top, width, height,
    *,
    accent: RGBColor = ACCENT_PURPLE,
    fill: RGBColor = SURFACE,
    fill_alpha: float = 0.55,
):
    """Glass-ish card with a thin coloured border."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.08
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    _set_shape_alpha(card, fill_alpha)
    card.line.color.rgb = accent
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    return card


def add_chip(slide, left, top, text: str, color: RGBColor, *, width: float = 1.4, height: float = 0.32):
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(width), Inches(height)
    )
    chip.adjustments[0] = 0.5
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = SNAKE_WHITE


# ----------------------------------------------------------------------------
# Each slide
# ----------------------------------------------------------------------------
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_dark_bg(slide, prs)

    add_text(
        slide, Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
        "Network Anomaly Detector",
        size=54, bold=True, color=TEXT, align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.7), Inches(3.2), Inches(12), Inches(0.6),
        "ML-powered traffic analysis · auto-block in real time",
        size=22, color=ACCENT_PURPLE, align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
        "Personal Research Project · Cybersecurity Attack & Defend",
        size=16, color=TEXT_MUTED,
    )
    add_text(
        slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
        "Angel Rusev  ·  Fontys University of Applied Sciences  ·  Spring 2026",
        size=13, color=TEXT_MUTED,
    )


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "The problem", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "Traditional firewalls only catch what they've seen before.",
             size=20, color=ACCENT_PURPLE)

    add_bullets(
        slide, Inches(0.7), Inches(2.4), Inches(12), Inches(4),
        [
            "Networks generate millions of connections per day — too much for humans.",
            "Signature-based tools (Snort, Suricata) miss anything new or unusual.",
            "Analysts get flooded with alerts and start ignoring them.",
            "We need something that learns what normal looks like and flags the rest.",
        ],
        size=20,
    )


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "Our solution", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "A small machine-learning model that reads each network flow and decides.",
             size=20, color=ACCENT_PURPLE)

    add_bullets(
        slide, Inches(0.7), Inches(2.4), Inches(12), Inches(4),
        [
            "Watches packets going across the network in real time.",
            "Tells you what kind of traffic it is — normal or one of 7 attack types.",
            "Blocks the attacker's IP automatically when it's confident.",
            "Shows everything in a dashboard so you can see it happening.",
        ],
        size=20,
    )


def slide_what_we_detect(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "What we detect", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "Eight categories — one normal, seven attacks.",
             size=20, color=ACCENT_PURPLE)

    # Grid of chips
    rows = [
        ("BENIGN",       "Normal browsing, API calls"),
        ("PortScan",     "Mapping open ports (nmap)"),
        ("DDoS",         "Flood from many sources"),
        ("DoS",          "Flood from one source"),
        ("BruteForce",   "Guessing passwords (hydra)"),
        ("WebAttack",    "SQL injection, XSS (sqlmap)"),
        ("Botnet",       "Infected machine calling home"),
        ("Infiltration", "Quietly stealing data"),
    ]
    cols = 2
    chip_w = 1.6
    desc_w = 4.5
    start_x = 0.9
    start_y = 2.4
    row_h = 0.65
    for i, (cat, desc) in enumerate(rows):
        col = i % cols
        row = i // cols
        x = start_x + col * (chip_w + desc_w + 0.4)
        y = start_y + row * row_h
        add_chip(slide, Inches(x), Inches(y), cat, CAT_COLORS[cat], width=chip_w, height=0.36)
        add_text(slide, Inches(x + chip_w + 0.15), Inches(y + 0.02),
                 Inches(desc_w), Inches(0.4), desc,
                 size=16, color=TEXT_MUTED)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "How it fits together", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "Three small services, each does one job.",
             size=20, color=ACCENT_PURPLE)

    # Three cards in a row
    card_y = Inches(2.5)
    card_h = Inches(3.3)
    card_w = Inches(3.9)
    gap = Inches(0.3)
    start_x = Inches(0.7)

    cards = [
        ("ATTACKER", "Generates real attack traffic\n(nmap, hping3, hydra, sqlmap)", ACCENT_INDIGO),
        ("DEFENDER", "Sees the packets\nClassifies them\nBlocks the IP", ACCENT_PURPLE),
        ("BRAIN (API)", "Holds the ML model\nReturns a verdict\nRecords alerts", SEV_INFO),
    ]
    for i, (title, body, accent) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, card_y, card_w, card_h, accent=accent, fill_alpha=0.5)
        add_text(slide, x + Inches(0.25), card_y + Inches(0.25),
                 card_w - Inches(0.5), Inches(0.5),
                 title, size=20, bold=True, color=accent)
        add_text(slide, x + Inches(0.25), card_y + Inches(0.95),
                 card_w - Inches(0.5), card_h - Inches(1.0),
                 body, size=16, color=TEXT)

    # Arrows under the cards: attacker → defender → brain (curved)
    arrow_y = card_y + card_h + Inches(0.25)
    add_text(slide, Inches(0.7), arrow_y, Inches(12), Inches(0.4),
             "attacker  →  defender  ←→  brain",
             size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER, font=FONT_MONO)


def slide_brain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "The brain — what the ML model learned", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "We trained on a public dataset of real attacks.",
             size=20, color=ACCENT_PURPLE)

    add_bullets(
        slide, Inches(0.7), Inches(2.4), Inches(7), Inches(4),
        [
            "Dataset: CICIDS2017 — millions of labelled network flows.",
            "Algorithm: Gradient Boosting (a tree-based classifier).",
            "Looks at ~20 features per flow: packet sizes, timing, flags.",
            "Outputs: the most likely category + a confidence score.",
        ],
        size=18,
    )

    # A "model card" on the right
    card_x = Inches(8.5)
    card_y = Inches(2.4)
    card_w = Inches(4.2)
    card_h = Inches(3.3)
    add_card(slide, card_x, card_y, card_w, card_h, accent=ACCENT_PURPLE, fill_alpha=0.55)
    add_text(slide, card_x + Inches(0.25), card_y + Inches(0.25),
             card_w - Inches(0.5), Inches(0.4),
             "MODEL", size=12, bold=True, color=TEXT_MUTED)
    add_text(slide, card_x + Inches(0.25), card_y + Inches(0.7),
             card_w - Inches(0.5), Inches(0.5),
             "gradient_boosting", size=22, bold=True, color=ACCENT_PURPLE,
             font=FONT_MONO)
    add_text(slide, card_x + Inches(0.25), card_y + Inches(1.4),
             card_w - Inches(0.5), Inches(0.4),
             "INPUT", size=12, bold=True, color=TEXT_MUTED)
    add_text(slide, card_x + Inches(0.25), card_y + Inches(1.85),
             card_w - Inches(0.5), Inches(0.4),
             "20 flow features", size=18, color=TEXT)
    add_text(slide, card_x + Inches(0.25), card_y + Inches(2.35),
             card_w - Inches(0.5), Inches(0.4),
             "OUTPUT", size=12, bold=True, color=TEXT_MUTED)
    add_text(slide, card_x + Inches(0.25), card_y + Inches(2.8),
             card_w - Inches(0.5), Inches(0.4),
             "verdict + score 0–1", size=18, color=TEXT)


def slide_detector(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "The defender — sees and reacts", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "Steady loop, four steps, no human in the loop.",
             size=20, color=ACCENT_PURPLE)

    steps = [
        ("1  CAPTURE",   "tcpdump on the network card",  ACCENT_INDIGO),
        ("2  EXTRACT",   "CICFlowMeter groups packets into flows",  ACCENT_PURPLE),
        ("3  CLASSIFY",  "Asks the brain: is this an attack?",  SEV_LOW),
        ("4  ENFORCE",   "If yes — iptables blocks the source IP",  SEV_CRIT),
    ]
    y = 2.4
    for i, (title, desc, color) in enumerate(steps):
        add_chip(slide, Inches(0.7), Inches(y + i*0.8), title, color, width=1.7, height=0.5)
        add_text(slide, Inches(2.7), Inches(y + i*0.8 + 0.08), Inches(10), Inches(0.5),
                 desc, size=18, color=TEXT)


def slide_attacker(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "The attacker — real tools from the course", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "Each preset in the dashboard maps to one real tool.",
             size=20, color=ACCENT_PURPLE)

    rows = [
        ("PortScan",    "nmap -sS",        CAT_COLORS["PortScan"]),
        ("DDoS / DoS",  "hping3 --flood",  CAT_COLORS["DDoS"]),
        ("BruteForce",  "hydra ssh://...", CAT_COLORS["BruteForce"]),
        ("WebAttack",   "sqlmap -u ...",   CAT_COLORS["WebAttack"]),
        ("Botnet",      "scripted beacon", CAT_COLORS["Botnet"]),
        ("Infiltration","large uploads",   CAT_COLORS["Infiltration"]),
    ]
    y = 2.4
    for i, (preset, tool, color) in enumerate(rows):
        add_chip(slide, Inches(0.9), Inches(y + i*0.55), preset, color, width=1.7, height=0.4)
        add_text(slide, Inches(2.9), Inches(y + i*0.55 + 0.04),
                 Inches(9), Inches(0.4),
                 tool, size=17, color=TEXT, font=FONT_MONO)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "Does it actually work?", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "We ran a real nmap scan against the lab and checked the verdicts.",
             size=20, color=ACCENT_PURPLE)

    # Big number cards
    card_y = Inches(2.6)
    card_h = Inches(2.0)
    card_w = Inches(3.9)
    gap = Inches(0.25)
    start_x = Inches(0.7)

    metrics = [
        ("4 of 5", "flows classified as PortScan",                 ACCENT_PURPLE),
        ("1.00",   "confidence on the correct class",              SEV_INFO),
        ("0",      "false alarms on benign traffic",               SEV_LOW),
    ]
    for i, (big, small, color) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, card_y, card_w, card_h, accent=color, fill_alpha=0.5)
        add_text(slide, x + Inches(0.25), card_y + Inches(0.25),
                 card_w - Inches(0.5), Inches(0.9),
                 big, size=46, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT_MONO)
        add_text(slide, x + Inches(0.25), card_y + Inches(1.25),
                 card_w - Inches(0.5), Inches(0.6),
                 small, size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.7), Inches(5.0), Inches(12), Inches(0.6),
             "The fifth flow targeted an open port — the model called it DDoS instead. "
             "Honest borderline case, worth noting.",
             size=14, color=TEXT_MUTED)


def slide_lab_setup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "How it all runs", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "Everything in Docker — reproducible, no setup pain.",
             size=20, color=ACCENT_PURPLE)

    add_bullets(
        slide, Inches(0.7), Inches(2.4), Inches(12), Inches(4),
        [
            "API + model: Python + FastAPI + scikit-learn.",
            "Frontend: React + TypeScript + Plotly charts.",
            "Lab: three containers — attacker, defender, brain — talking over a Docker bridge.",
            "One command: docker compose up — and the whole demo is live.",
        ],
        size=20,
    )


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(1.2),
             "Live demo", size=72, bold=True, color=ACCENT_PURPLE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.7), Inches(3.0), Inches(12), Inches(0.6),
             "Switch to the dashboard.", size=24, color=TEXT,
             align=PP_ALIGN.CENTER)

    add_bullets(
        slide, Inches(2.5), Inches(4.2), Inches(8.5), Inches(2.5),
        [
            "Pick a preset (PortScan, DDoS, BruteForce…).",
            "Click \"Run on lab\" — real tool, real packets.",
            "Watch the verdict appear in the Alerts tab.",
            "Show the iptables rule blocking the attacker.",
            "Run it again — second attack times out.",
        ],
        size=18,
    )


def slide_next(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.8),
             "What's next", size=36, bold=True, color=TEXT)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
             "Beyond the PRP — integrating into the IAM project.",
             size=20, color=ACCENT_PURPLE)

    add_bullets(
        slide, Inches(0.7), Inches(2.4), Inches(12), Inches(4),
        [
            "Plug the detector into the CYBERGROUP IAM project (Keycloak-protected APIs).",
            "Retrain on traffic we capture ourselves — closer to what we'd see in production.",
            "Add an analyst-confirm step before permanent blocks (no surprise lockouts).",
            "Move from one defender container to an inline gateway pattern.",
        ],
        size=18,
    )


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    add_text(slide, Inches(0.7), Inches(2.8), Inches(12), Inches(1.5),
             "Thank you", size=72, bold=True, color=TEXT,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.7), Inches(4.3), Inches(12), Inches(0.5),
             "Questions?", size=28, color=ACCENT_PURPLE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
             "Angel Rusev  ·  Fontys University of Applied Sciences  ·  Spring 2026",
             size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ----------------------------------------------------------------------------
# Build
# ----------------------------------------------------------------------------
def main() -> None:
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_what_we_detect(prs)
    slide_architecture(prs)
    slide_brain(prs)
    slide_detector(prs)
    slide_attacker(prs)
    slide_lab_setup(prs)
    slide_results(prs)
    slide_demo(prs)
    slide_next(prs)
    slide_thanks(prs)

    out = Path(__file__).resolve().parent / "AnomalyDetector-Presentation.pptx"
    prs.save(out)
    print(f"saved: {out}  ({out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
